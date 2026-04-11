"""Document parser.

Supports: PDF, DOCX, PPTX, XLSX, HTML, Markdown, plain text.
Primary: Docling for structured parsing.
Fallbacks: pypdf, python-docx, python-pptx, openpyxl for robustness on CPU VPS.
"""
from __future__ import annotations

import io
import logging
import tempfile
from dataclasses import dataclass
from pathlib import Path

log = logging.getLogger(__name__)

# All MIME types we can handle
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
DOC_MIME = "application/msword"
PPT_MIME = "application/vnd.ms-powerpoint"
XLS_MIME = "application/vnd.ms-excel"


@dataclass
class ParsedSection:
    order: int
    page: int | None
    section_title: str | None
    text: str


def parse_document(content: bytes, mime_type: str, filename: str) -> list[ParsedSection]:
    ext = Path(filename).suffix.lower()

    # --- PDF ---
    if mime_type == "application/pdf" or ext == ".pdf":
        return _parse_pdf(content, filename)

    # --- DOCX / DOC ---
    if mime_type in {DOCX_MIME, DOC_MIME} or ext in {".docx", ".doc"}:
        return _parse_docx(content)

    # --- PPTX / PPT ---
    if mime_type in {PPTX_MIME, PPT_MIME} or ext in {".pptx", ".ppt"}:
        return _parse_pptx(content)

    # --- XLSX / XLS ---
    if mime_type in {XLSX_MIME, XLS_MIME} or ext in {".xlsx", ".xls", ".csv"}:
        return _parse_xlsx(content, ext)

    # --- HTML ---
    if mime_type == "text/html" or ext == ".html":
        return _parse_html(content, filename)

    # --- Text / Markdown / JSON / XML / CSV ---
    if mime_type.startswith("text/") or mime_type in {"application/json", "application/xml"}:
        text = content.decode("utf-8", errors="replace")
        return _split_markdown(text)

    # Last-resort: decode with replacement.
    text = content.decode("utf-8", errors="replace")
    return _split_markdown(text)


# ---------------------------------------------------------------------------
# Format-specific parsers
# ---------------------------------------------------------------------------


def _parse_pdf(content: bytes, filename: str) -> list[ParsedSection]:
    """Parse PDF: try Docling first, fallback to pypdf."""
    # Try Docling
    try:
        return _docling_parse(content, filename)
    except Exception as exc:
        log.warning("docling_pdf_failed: %s – trying pypdf", exc)

    # Fallback: pypdf
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        pages: list[ParsedSection] = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(ParsedSection(order=i, page=i + 1, section_title=None, text=text.strip()))
        if pages:
            log.info("pypdf_ok: %d pages", len(pages))
            return pages
    except Exception as pdf_exc:
        log.warning("pypdf_failed: %s", pdf_exc)

    raise RuntimeError("PDF parsing failed with all methods")


def _parse_docx(content: bytes) -> list[ParsedSection]:
    """Parse DOCX using python-docx."""
    try:
        from docx import Document  # type: ignore

        doc = Document(io.BytesIO(content))
        sections: list[ParsedSection] = []
        order = 0
        current_title: str | None = None
        current_lines: list[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Detect headings
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                if current_lines:
                    sections.append(ParsedSection(
                        order=order, page=None, section_title=current_title,
                        text="\n".join(current_lines),
                    ))
                    order += 1
                    current_lines = []
                current_title = text
            else:
                current_lines.append(text)

        if current_lines:
            sections.append(ParsedSection(
                order=order, page=None, section_title=current_title,
                text="\n".join(current_lines),
            ))

        # Also extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                order += 1
                sections.append(ParsedSection(
                    order=order, page=None, section_title="Tableau",
                    text="\n".join(rows),
                ))

        if sections:
            log.info("docx_ok: %d sections", len(sections))
            return sections
    except Exception as exc:
        log.warning("docx_parse_failed: %s", exc)

    raise RuntimeError("DOCX parsing failed")


def _parse_pptx(content: bytes) -> list[ParsedSection]:
    """Parse PPTX using python-pptx."""
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(io.BytesIO(content))
        sections: list[ParsedSection] = []

        for i, slide in enumerate(prs.slides):
            texts: list[str] = []
            title: str | None = None
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
            if texts:
                sections.append(ParsedSection(
                    order=i, page=i + 1, section_title=title or f"Diapositive {i + 1}",
                    text="\n".join(texts),
                ))

        if sections:
            log.info("pptx_ok: %d slides", len(sections))
            return sections
    except Exception as exc:
        log.warning("pptx_parse_failed: %s", exc)

    raise RuntimeError("PPTX parsing failed")


def _parse_xlsx(content: bytes, ext: str) -> list[ParsedSection]:
    """Parse XLSX/CSV using openpyxl or csv."""
    if ext == ".csv":
        import csv

        text = content.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        rows = [" | ".join(row) for row in reader]
        if rows:
            return [ParsedSection(order=0, page=None, section_title="Tableau CSV", text="\n".join(rows))]
        return [ParsedSection(order=0, page=None, section_title=None, text=text)]

    try:
        from openpyxl import load_workbook  # type: ignore

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sections: list[ParsedSection] = []

        for i, sheet in enumerate(wb.sheetnames):
            ws = wb[sheet]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = " | ".join(cells).strip()
                if line and line != "| " * len(cells):
                    rows.append(line)
            if rows:
                sections.append(ParsedSection(
                    order=i, page=None, section_title=f"Feuille: {sheet}",
                    text="\n".join(rows),
                ))

        wb.close()
        if sections:
            log.info("xlsx_ok: %d sheets", len(sections))
            return sections
    except Exception as exc:
        log.warning("xlsx_parse_failed: %s", exc)

    raise RuntimeError("XLSX parsing failed")


def _parse_html(content: bytes, filename: str) -> list[ParsedSection]:
    """Parse HTML: try Docling, fallback to basic tag stripping."""
    try:
        return _docling_parse(content, filename)
    except Exception:
        pass

    import re

    text = content.decode("utf-8", errors="replace")
    text = re.sub(r"<script[^>]*>.*?</script>", "", text, flags=re.DOTALL)
    text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    if text:
        return [ParsedSection(order=0, page=None, section_title=None, text=text)]
    raise RuntimeError("HTML parsing failed")


def _docling_parse(content: bytes, filename: str) -> list[ParsedSection]:
    """Try Docling structured parsing."""
    from docling.document_converter import DocumentConverter  # type: ignore

    converter = DocumentConverter()
    try:
        from docling.datamodel.base_models import DocumentStream  # type: ignore

        source = DocumentStream(name=filename, stream=io.BytesIO(content))
        result = converter.convert(source)
    except ImportError:
        suffix = Path(filename).suffix or ".bin"
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(content)
            tmp.flush()
            tmp_path = tmp.name
        try:
            result = converter.convert(source=tmp_path)
        finally:
            Path(tmp_path).unlink(missing_ok=True)

    md = result.document.export_to_markdown()
    return _split_markdown(md)


def _split_markdown(text: str) -> list[ParsedSection]:
    """Split on markdown headings into coarse sections."""
    text = text.replace("\x00", "")
    sections: list[ParsedSection] = []
    current_title: str | None = None
    current_lines: list[str] = []
    order = 0

    def flush() -> None:
        nonlocal order
        body = "\n".join(current_lines).strip()
        if body:
            sections.append(
                ParsedSection(order=order, page=None, section_title=current_title, text=body)
            )
            order += 1

    for line in text.splitlines():
        if line.startswith("#"):
            flush()
            current_title = line.lstrip("# ").strip() or None
            current_lines = []
        else:
            current_lines.append(line)
    flush()

    if not sections:
        sections.append(ParsedSection(order=0, page=None, section_title=None, text=text.strip()))
    return sections
